import logging
from io import BytesIO
from pathlib import Path
from typing import Set, Union

from docling_core.types.experimental import (
    BasePictureData,
    BaseTableData,
    DescriptionItem,
    DocItemLabel,
    DoclingDocument,
    GroupLabel,
    ImageRef,
    PictureItem,
    ProvenanceItem,
    TableCell,
    TableItem,
)
from docling_core.types.experimental.base import BoundingBox, CoordOrigin, Size
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE, PP_PLACEHOLDER
from pptx.util import Inches

from docling.backend.abstract_backend import DeclarativeDocumentBackend
from docling.datamodel.base_models import InputFormat

_log = logging.getLogger(__name__)


class MsPowerpointDocumentBackend(DeclarativeDocumentBackend):
    def __init__(self, path_or_stream: Union[BytesIO, Path], document_hash: str):
        super().__init__(path_or_stream, document_hash)
        self.namespaces = {
            "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
            "c": "http://schemas.openxmlformats.org/drawingml/2006/chart",
            "p": "http://schemas.openxmlformats.org/presentationml/2006/main",
        }
        # Powerpoint file:
        self.path_or_stream = path_or_stream
        return

    def is_valid(self) -> bool:
        return True

    def is_paginated(cls) -> bool:
        False

    def unload(self):
        if isinstance(self.path_or_stream, BytesIO):
            self.path_or_stream.close()

        self.path_or_stream = None

    @classmethod
    def supported_formats(cls) -> Set[InputFormat]:
        return {InputFormat.PPTX}

    def convert(self) -> DoclingDocument:
        # Parses the PPTX into a structured document model.
        doc = DoclingDocument(description=DescriptionItem(), name="dummy")
        pptx_obj = None
        try:
            pptx_obj = Presentation(self.path_or_stream)
        except Exception:
            _log.error("could not parse pptx")
            return doc
        doc = self.walk_linear(pptx_obj, doc)

        return doc

    def generate_prov(self, shape, slide_ind, text=""):
        left = shape.left
        top = shape.top
        width = shape.width
        height = shape.height
        shape_bbox = [left, top, left + width, top + height]
        shape_bbox = BoundingBox.from_tuple(shape_bbox, origin=CoordOrigin.BOTTOMLEFT)
        # prov = [{"bbox": shape_bbox, "page": parent_slide, "span": [0, len(text)]}]
        prov = ProvenanceItem(
            page_no=slide_ind, charspan=[0, len(text)], bbox=shape_bbox
        )

        return prov

    def handle_text_elements(self, shape, parent_slide, slide_ind, doc):
        is_a_list = False
        for paragraph in shape.text_frame.paragraphs:
            bullet_type = "None"
            # Check if paragraph is a bullet point using the `element` XML
            p = paragraph._element
            if (
                p.find(".//a:buChar", namespaces={"a": self.namespaces["a"]})
                is not None
            ):
                bullet_type = "Bullet"
                is_a_list = True
            elif (
                p.find(".//a:buAutoNum", namespaces={"a": self.namespaces["a"]})
                is not None
            ):
                bullet_type = "Numbered"
                is_a_list = True
            else:
                is_a_list = False
            if paragraph.level > 0:
                # Most likely a sub-list
                is_a_list = True
            list_text = paragraph.text.strip()

            prov = self.generate_prov(shape, slide_ind, shape.text.strip())

            if is_a_list:
                new_list = doc.add_group(
                    label=GroupLabel.LIST, name=f"list", parent=parent_slide
                )
            else:
                new_list = None

            # for element in p.iter():
            for e in p.iterfind(".//a:r", namespaces={"a": self.namespaces["a"]}):
                if len(e.text.strip()) > 0:
                    if (
                        p.find(".//a:buChar", namespaces={"a": self.namespaces["a"]})
                        is not None
                    ):
                        bullet_type = "Bullet"
                        e_is_a_list_item = True
                    elif (
                        p.find(".//a:buAutoNum", namespaces={"a": self.namespaces["a"]})
                        is not None
                    ):
                        bullet_type = "Numbered"
                        e_is_a_list_item = True
                    else:
                        e_is_a_list_item = False
                    if e_is_a_list_item:
                        doc.add_text(
                            label=DocItemLabel.LIST_ITEM,
                            parent=new_list,
                            text=e.text.strip(),
                            prov=prov,
                        )
                    else:
                        doc.add_text(
                            label=DocItemLabel.PARAGRAPH,
                            parent=parent_slide,
                            text=e.text.strip(),
                            prov=prov,
                        )
        return

    def handle_title(self, shape, parent_slide, slide_ind, doc):
        placeholder_type = shape.placeholder_format.type
        txt = shape.text.strip()
        prov = self.generate_prov(shape, slide_ind, txt)

        if len(txt.strip()) > 0:
            # title = slide.shapes.title.text if slide.shapes.title else "No title"
            if placeholder_type in [PP_PLACEHOLDER.CENTER_TITLE, PP_PLACEHOLDER.TITLE]:
                _log.info(f"Title found: {shape.text}")
                doc.add_text(
                    label=DocItemLabel.TITLE, parent=parent_slide, text=txt, prov=prov
                )
            elif placeholder_type == PP_PLACEHOLDER.SUBTITLE:
                _log.info(f"Subtitle found: {shape.text}")
                # Using DocItemLabel.FOOTNOTE, while SUBTITLE label is not avail.
                doc.add_text(
                    label=DocItemLabel.SECTION_HEADER,
                    parent=parent_slide,
                    text=txt,
                    prov=prov,
                )
        return

    def handle_pictures(self, shape, parent_slide, slide_ind, doc):
        # shape has picture
        prov = self.generate_prov(shape, slide_ind, "")
        doc.add_picture(
            data=BasePictureData(), parent=parent_slide, caption=None, prov=prov
        )
        return

    def handle_tables(self, shape, parent_slide, slide_ind, doc):
        # Handling tables, images, charts
        if shape.has_table:
            table = shape.table
            table_xml = shape._element

            prov = self.generate_prov(shape, slide_ind, "")

            num_cols = 0
            num_rows = len(table.rows)
            tcells = []
            # Access the XML element for the shape that contains the table
            table_xml = shape._element

            for row_idx, row in enumerate(table.rows):
                if len(row.cells) > num_cols:
                    num_cols = len(row.cells)
                for col_idx, cell in enumerate(row.cells):
                    # Access the XML of the cell (this is the 'tc' element in table XML)
                    cell_xml = table_xml.xpath(
                        f".//a:tbl/a:tr[{row_idx + 1}]/a:tc[{col_idx + 1}]"
                    )

                    if not cell_xml:
                        continue  # If no cell XML is found, skip

                    cell_xml = cell_xml[0]  # Get the first matching XML node
                    row_span = cell_xml.get("rowSpan")  # Vertical span
                    col_span = cell_xml.get("gridSpan")  # Horizontal span

                    if row_span is None:
                        row_span = 1
                    else:
                        row_span = int(row_span)

                    if col_span is None:
                        col_span = 1
                    else:
                        col_span = int(col_span)

                    icell = TableCell(
                        text=cell.text.strip(),
                        row_span=row_span,
                        col_span=col_span,
                        start_row_offset_idx=row_idx,
                        end_row_offset_idx=row_idx + row_span,
                        start_col_offset_idx=col_idx,
                        end_col_offset_idx=col_idx + col_span,
                        col_header=False,
                        row_header=False,
                    )
                    if len(cell.text.strip()) > 0:
                        tcells.append(icell)
            # Initialize Docling BaseTableData
            data = BaseTableData(num_rows=num_rows, num_cols=num_cols, table_cells=[])
            # Populate
            for tcell in tcells:
                data.table_cells.append(tcell)
            if len(tcells) > 0:
                # If table is not fully empty...
                # Create Docling table
                doc.add_table(data=data, prov=prov)
        return

    def walk_linear(self, pptx_obj, doc) -> DoclingDocument:
        # Units of size in PPTX by default are EMU units (English Metric Units)
        slide_width = pptx_obj.slide_width
        slide_height = pptx_obj.slide_height

        text_content = []

        max_levels = 10
        parents = {}
        for i in range(0, max_levels):
            parents[i] = None

        # Loop through each slide
        for slide_num, slide in enumerate(pptx_obj.slides):
            slide_ind = pptx_obj.slides.index(slide)
            parent_slide = doc.add_group(
                name=f"slide-{slide_ind}", label=GroupLabel.CHAPTER, parent=parents[0]
            )

            size = Size(width=slide_width, height=slide_height)
            parent_page = doc.add_page(page_no=slide_ind, size=size)
            # parent_page = doc.add_page(page_no=slide_ind, size=size, hash=hash)

            # Loop through each shape in the slide
            for shape in slide.shapes:

                if shape.has_table:
                    # Handle Tables
                    self.handle_tables(shape, parent_slide, slide_ind, doc)

                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    # Handle Tables
                    self.handle_pictures(shape, parent_slide, slide_ind, doc)

                # If shape doesn't have any text, move on to the next shape
                if not hasattr(shape, "text"):
                    continue
                if shape.text is None:
                    continue
                if len(shape.text.strip()) == 0:
                    continue
                if not shape.has_text_frame:
                    _log.warn("Warning: shape has text but not text_frame")
                    continue

                if shape.is_placeholder:
                    # Handle Titles (Headers) and Subtitles
                    # Check if the shape is a placeholder (titles are placeholders)
                    self.handle_title(shape, parent_slide, slide_ind, doc)
                else:
                    # Handle other text elements, including lists (bullet lists, numbered lists)
                    self.handle_text_elements(shape, parent_slide, slide_ind, doc)
                # figures...
                # doc.add_figure(data=BaseFigureData(), parent=self.parents[self.level], caption=None)

        return doc